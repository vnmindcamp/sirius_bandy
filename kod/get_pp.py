from pptx import Presentation
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import general_functions as gf
from get_stats import Stats
from get_data import Game
from get_plot import Plot
from compile_stats import CompileStats
import os
import constants
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
# importeras typ bara för läslighet och typing 
import pptx.shapes
import pptx.table  
import pptx.slide 

# för att chilla med borders
from pptx.oxml.xmlchemy import OxmlElement



class PP:
    # class variables
    # sirius specific colors
    text_color = RGBColor(0, 0, 0)

    text_color_main = RGBColor(255, 0, 0)
    text_color_opponent = RGBColor(0, 0, 0)
    image_color_main = RGBColor(255, 0, 0)
    image_color_opponent = RGBColor(0, 0, 0)
    background_color = RGBColor(255,255,255)  
    background_image = "..\\..\\bilder\\pptx\\faded_background.png" #"..\\..\\bilder\\pptx\\background.png"

    # relative link to image folder
    logo_image_link = "..\\..\\bilder\\logos\\"
    auto_image_link = "..\\..\\bilder\\autogen\\"
    player_image_link = "..\\..\\bilder\\player_pics\\"   

    # contructor
    def __init__(self, stats: Stats, N: int = 1) -> None:
        # we will always print our team first
        self.ordered_teams = [stats.main_team, stats.opposite_team(stats.main_team)]
        self.stats = stats
        self.plot = Plot(stats, N = N)
        self.pres = Presentation()
        # måste kalla make_game_report
        #self.make_game_report()
        return

    # POWER POINT REPORTS 
    def make_game_report(self) -> None:
        '''calls the methods needed to make a game report presentation '''
        self.make_game_report_front_page()
        self.make_game_report_overview_stats_page()
        self.make_game_report_duels_page()
        #self.make_game_report_scimmages_page()
        self.make_game_report_before_and_after_table_page()
        self.make_game_report_shot_types_page(show_xg = True)
        self.make_single_image_page(self.plot.make_expected_goals_over_time_image(main_team_color=gf.rgb255_to_rgb1(PP.image_color_main), other_team_color=gf.rgb255_to_rgb1(PP.image_color_opponent)), 'Mål och XG\n', from_top = 0)
        self.make_game_report_slot_page()
        #self.make_game_report_shot_origins_page()
        self.make_game_report_goals_stats_page()
        self.make_single_image_page(self.plot.make_all_duels_locations_image(number_text=True), 'Alla närkamper och brytningar per zon')
        self.make_single_image_page(self.plot.make_duel_winners_per_locations_image(text_type='procent'), f"{gf.get_nickname(self.stats.main_team,'full')} vunna närkamper och brytningar per zon")
        # make_duel_zones_per_team_image
        self.make_single_image_page(self.plot.make_duel_zones_per_team_image(team_in_possession=self.stats.opposite_team(self.stats.main_team), number_text=True), 'Defensiva närkamper och brytningar per zon')
        self.make_single_image_page(self.plot.make_duel_winners_per_zone_and_team_image(team_in_possession=self.stats.opposite_team(self.stats.main_team), text_type='procent'), f"{gf.get_nickname(self.stats.main_team,'full')} vunna defensiva närkamper och brytningar per zon")
        self.make_single_image_page(self.plot.make_all_freeshot_locations_image(number_text=True), 'Frislag per zon')
        self.make_single_image_page(self.plot.make_freeshots_made_per_locations_image(text_type='procent'), f"{gf.get_nickname(self.stats.main_team,'full')}s frislag per zon")

        self.make_game_report_per_time_page()

        #fourty_image = self.plot.make_per_minute_bars(self.stats.prints['40'], ylabel='40-spel (antal)', color=PP.image_color_main)
        #self.make_single_image_page(fourty_image, title_text=f"Spridning av {constants.nicknames[self.stats.main_team]['full']} {sum(self.stats.prints['40'])} st 40-spel", from_top=0.3)
        #self.make_game_report_attacks_and_fourty_page()

        self.save_presentation()

    def make_season_report(self, filename: str = None) -> None:
        '''calls the methods needed to make a season report presentation'''
        if filename != None:
            self.stats.out = filename

        self.make_season_report_front_page()
        self.make_season_report_overview_stats_page()
        self.make_season_report_shots_page(show_xg=True)
        self.make_season_report_slot_page()
        self.make_season_report_duels_page()
        self.make_game_report_before_and_after_table_page()
        self.make_single_image_page(self.plot.make_all_duels_locations_image(number_text=True), 'Alla närkamper och brytningar per zon')
        self.make_single_image_page(self.plot.make_duel_winners_per_locations_image(text_type='procent'), f"{gf.get_nickname(self.stats.main_team,'full')} vunna närkamper och brytningar per zon")
        # defensiva närkamper
        self.make_single_image_page(self.plot.make_duel_zones_per_team_image(team_in_possession=self.stats.opposite_team(self.stats.main_team), number_text=True), f'{gf.get_nickname(self.stats.main_team,"full")} defensiva närkamper och brytningar per zon')
        self.make_single_image_page(self.plot.make_duel_winners_per_zone_and_team_image(team_in_possession=self.stats.opposite_team(self.stats.main_team), text_type='procent'), f"{gf.get_nickname(self.stats.main_team,'full')} vunna defensiva närkamper och brytningar per zon")        
        # offensiva närkamper
        self.make_single_image_page(self.plot.make_duel_zones_per_team_image(team_in_possession=self.stats.main_team, number_text=True), f'{gf.get_nickname(self.stats.main_team,"full")} offensiva närkamper och brytningar per zon')
        self.make_single_image_page(self.plot.make_duel_winners_per_zone_and_team_image(team_in_possession=self.stats.main_team, text_type='procent'), f"{gf.get_nickname(self.stats.main_team,'full')} vunna offensiva närkamper och brytningar per zon")        
        
        self.make_season_report_corners_page()
        self.make_season_report_long_shot_targets_page()
        self.make_season_report_possession_after_long_shots_page()
        self.make_season_report_after_long_shots_attacking_team_page()
        self.make_season_report_after_long_shots_defending_team_page()

        self.save_presentation()
    
    def make_comparative_report(self, other: Stats, filename: str) -> None: 
        '''calls the methods needed to make a comparative report presentation'''
        self.other = other
        self.stats.out = filename

        self.make_comparative_report_front_page()
        self.make_comparative_report_overview_stats_page()
        self.make_comparative_report_shots_for_page()
        self.make_comparative_report_shots_against_page()
        self.make_comparative_report_slot_for_page()
        self.make_comparative_report_slot_against_page()
        


        self.save_presentation()

    def make_player_report(self, players: list, filename: str = 'spelarrapport') -> None:
        '''creates a player report'''
        self.stats.out = filename
        self.make_player_report_front_page(players)
        for player in players:
            self.make_player_page(player)
        self.save_presentation()
        return 

    # static methods
    def iter_cells(table: pptx.table.Table) -> pptx.table._Cell:
        '''yeilds the cells of a table, 
            used to style them'''
        for row in table.rows:
            for cell in row.cells:
                yield cell

    def remove_elemet(element) -> None:
        '''removes element from PP'''
        # ta_bort = slide.placeholders[1]
        p = element._sp
        parent_element = p.getparent()
        parent_element.remove(p)

    # non-static methods

    def set_background_color(self, slide: pptx.slide.Slide) -> None:
        '''sets the background color of current slide to match
            class variable background_color'''
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PP.background_color
        #fill.fore_color.brightness = 0.1
        #fill.transparency = 0.1
        return   

    def set_background_image(self, slide: pptx.slide.Slide, img_path: str = None) -> None:
        '''sets the background image of current slide to match
            class variable background_image'''  
        left = top = Inches(0)
        if img_path == None:
            img_path = PP.background_image
        pic = slide.shapes.add_picture(img_path, left, top, width=self.pres.slide_width, height=self.pres.slide_height)
        # This moves it to the background
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    
    def get_team_text_color(self, team: str) -> RGBColor:
        '''returns the team color'''
        return PP.text_color_main if team == self.stats.main_team else PP.text_color_opponent

    def add_logo_images(self, slide, from_left = 0.7, from_top = 0.4, width = 2) -> None:
        '''adds the logos of the teams to the slide
            units in Inches''' 
        for i, team in enumerate(self.ordered_teams):
            img = PP.logo_image_link + gf.get_logo_image(team)
            if i != 0: 
                # converting to inches by dividing by 914400 ??????
                from_left = self.pres.slide_width/914400 - from_left - width 
            slide.shapes.add_picture(img, Inches(from_left), Inches(from_top), Inches(width))
    
    def add_main_team_logo(self, slide, from_left = 0.7, from_top = 0.4, width = 2) -> None:
        '''adds the logo of the main team to the slide
            units in Inches'''
        for i in range(2):
            img = PP.logo_image_link + gf.get_logo_image(self.stats.main_team) 
            if i != 0: 
                # converting to inches by dividing by 914400 ??????
                from_left = self.pres.slide_width/914400 - from_left - width 
            slide.shapes.add_picture(img, Inches(from_left), Inches(from_top), Inches(width))
    
    def add_player_photo(self, player: int, slide, from_right = 0.7, from_top = 0.4, width = 5) -> None:
        '''adds the photo of said player to the slide
            untis in inches'''
        img = PP.player_image_link + gf.get_player_info(player, 'image')
        
        # converting to inches by dividing by 914400 ??????
        from_right = self.pres.slide_width/914400 - from_right - width 
        slide.shapes.add_picture(img, Inches(from_right), Inches(from_top), Inches(width))
    
    def make_game_report_front_page(self) -> None:
        '''makes the front page layout'''
        slide_register = self.pres.slide_layouts[0]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_color(slide)
        title = slide.shapes.title
        title.text = " - ".join(gf.get_nickname(team, 'full') for team in self.ordered_teams) 
        title.text_frame.paragraphs[0].font.color.rgb = PP.text_color

        subtitle = slide.placeholders[1]
        subtitle.text = ' - '.join(str(sum(self.stats.prints['score'][team].values())) for team in self.ordered_teams) 
        subtitle.text_frame.paragraphs[0].font.color.rgb = PP.text_color
        self.add_logo_images(slide, width = 2)
        #self.set_background_color(slide)

    def make_player_report_front_page(self, players: list) -> None:
        '''makes the player page layout'''
        slide_register = self.pres.slide_layouts[0]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_image(slide)
        self.add_main_team_logo(slide, width=3)

        title = slide.shapes.title
        title.text = 'Spelarrapport' 
        title.text_frame.paragraphs[0].font.color.rgb = PP.text_color

        subtitle = slide.placeholders[1]
        subtitle.text = f'Rapport för {gf.readable_number(len(players))} spelare'
        # subtitle.text_frame.paragraphs[0].font.color.rgb = PP.text_color
   

    def make_season_report_front_page(self) -> None:
        '''makes the front page layout'''
        slide_register = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_color(slide)
        self.add_main_team_logo(slide)
        title = slide.shapes.title
        #title.top = Inches(2)
        title.text = f"Säsongsrapport för \n{gf.get_nickname(self.stats.main_team, 'full')}"
        title.text_frame.paragraphs[0].font.color.rgb = PP.text_color

    def make_comparative_report_front_page(self) -> None:
        '''makes the front page layout'''
        slide_register = self.pres.slide_layouts[0]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_color(slide)
        title = slide.shapes.title
        title.text = "Jämförelserapport"
        title.text_frame.paragraphs[0].font.color.rgb = PP.text_color

        subtitle = slide.placeholders[1]
        subtitle.text = f"{' - '.join(gf.get_nickname(team, 'full') for team in self.ordered_teams)} och {gf.readable_number(self.other.number_of_games)} halvlekar"
        subtitle.text_frame.paragraphs[0].font.color.rgb = PP.text_color
        self.add_main_team_logo(slide, width = 2)
        #self.set_background_color(slide)

    def make_player_page(self, player: int) -> None:
        '''creates a player page, assumes that the player appears in the constants.players'''
        slide_register = self.pres.slide_layouts[7]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_image(slide)
        #self.set_background_color(slide)
        self.add_player_photo(player, slide)
        title = slide.shapes.title 
        title.text = f"{player}: {gf.get_player_info(player, 'name')}"
        title.text_frame.paragraphs[0].font.size = Pt(27)

        PP.remove_elemet(slide.placeholders[1])

        player_text = slide.placeholders[2]
        player_stats = self.stats.get_player_stats_dict(player)
        rows_of_text = [
          #  f"{gf.get_player_info(player, 'position')}",
            f"Mål/XG: \n\t{player_stats['mål']}/{round(player_stats['xg'], 2)}",
        ]
        rows_of_text.append(f"Mål/skott: \n\tTotalt: {player_stats['mål']}/{player_stats['skott']}")
        for st in player_stats['skottyp']:
            goals = 0 if st not in player_stats['målskottyper'] else player_stats['målskottyper'][st]
            shots = player_stats['skottyp'][st]
            rows_of_text.append(f"\t{st}: {goals}/{shots}")
        
        rows_of_text.append(f"Passningar:")   
        if len(player_stats['passning']) > 0:    
            for p in player_stats['passning']:
                rows_of_text.append(f"\t{p}: {player_stats['passning'][p]}")
        else:
            rows_of_text.append(f"\t0")
        
        rows_of_text.append(f"Fixade hörnor:")
        rows_of_text.append(f"\t{player_stats['hörna']}")

        player_text.text = "\n".join(rows_of_text)
        player_text.text_frame.paragraphs[0].font.color.rgb = PP.text_color

    def make_comparative_report_overview_stats_page(self) -> None:
        '''makes the overview stats page layout'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        #self.set_background_color(slide)
        self.add_main_team_logo(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Övergripande \nhalvleksvis'

        # vänster sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[1] # first 1, then 3  
        bp1.text = f"{gf.get_nickname(self.other.main_team, 'short')} {self.other.number_of_games} halvlekar"
        #bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) #constants.colors[team][0]

        bp2 = bpb.placeholders[2]  # first 2, then 4
        res = bp2.text_frame.add_paragraph()
        res.text = f"Målsnitt för/emot: \n\t{round(sum(self.other.prints['goal types'][self.other.main_team].values()) / self.other.number_of_games, 1)} / {round(sum(self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)].values()) / self.other.number_of_games, 1)}"
        res.level = 0

        res = bp2.text_frame.add_paragraph()
        res.text = f"Skottsnitt för/emot: \n\t{round(self.other.prints['shot attempts'][self.other.main_team]  / self.other.number_of_games, 1)} / {round(self.other.prints['shot attempts'][self.other.opposite_team(self.other.main_team)]  / self.other.number_of_games, 1)}"
        res.level = 0 

        res = bp2.text_frame.add_paragraph()
        res.text = f"Hörnsnitt antal för/emot: \n\t{round(sum(self.other.prints['corners'][self.other.main_team].values()) / self.other.number_of_games, 1)} / {round(sum(self.other.prints['corners'][self.other.opposite_team(self.other.main_team)].values()) / self.other.number_of_games, 1)}"
        res.level = 0 

        res = bp2.text_frame.add_paragraph()
        res.text = f"Närkampssnitt för/emot: \n\t{round((self.other.prints['scrimmages'][self.other.main_team] + self.other.prints['interceptions'][self.other.main_team]) / self.other.number_of_games, 1)} / {round((self.other.prints['scrimmages'][self.other.opposite_team(self.other.main_team)] + self.other.prints['interceptions'][self.other.opposite_team(self.other.main_team)]) / self.other.number_of_games, 1)}"
        res.level = 0 
    
        res = bp2.text_frame.add_paragraph()
        res.text = f"Bollinnehav för/emot: \n\t{gf.sec_to_readable(gf.readable_to_sec(self.other.prints['possession'][self.other.main_team]) // self.other.number_of_games)} / {gf.sec_to_readable(gf.readable_to_sec(self.other.prints['possession'][self.other.opposite_team(self.other.main_team)]) // self.other.number_of_games)}"
        res.level = 0 


        #res = bp2.text_frame.add_paragraph()
        #res.text = f"Skottsnitt FÖR på mål/skottförsök: \n\t{round(self.other.prints['shots on goal'][self.other.main_team]  / self.other.number_of_games, 1)} / {round(self.other.prints['shot attempts'][self.other.main_team]  / self.other.number_of_games, 1)}"
        #res.level = 0


        # höger sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[3] # first 1, then 3  
        bp1.text = ' - '.join(gf.get_nickname(team, 'short') for team in self.ordered_teams)
        #bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) #constants.colors[team][0]

        bp2 = bpb.placeholders[4]  # first 2, then 4
        res = bp2.text_frame.add_paragraph()
        res.text = f"Mål för/emot: \n\t{sum(self.stats.prints['score'][self.stats.main_team].values())} / {sum(self.stats.prints['score'][self.stats.opposite_team(self.stats.main_team)].values())}"
        res.level = 0

        res = bp2.text_frame.add_paragraph()
        res.text = f"Skottförsök för/emot: \n\t{sum(self.stats.prints['shot types'][self.stats.main_team].values())} / {sum(self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)].values())}"
        res.level = 0

        res = bp2.text_frame.add_paragraph()
        res.text = f"Hörnsnitt för/emot: \n\t{sum(self.stats.prints['corners'][self.stats.main_team].values())} / {sum(self.stats.prints['corners'][self.stats.opposite_team(self.stats.main_team)].values())}"
        res.level = 0 

        res = bp2.text_frame.add_paragraph()
        res.text = f"Närkampssnitt för/emot: \n\t{self.stats.prints['scrimmages'][self.stats.main_team] + self.stats.prints['interceptions'][self.stats.main_team]} / {self.stats.prints['scrimmages'][self.stats.opposite_team(self.stats.main_team)] + self.stats.prints['interceptions'][self.stats.opposite_team(self.stats.main_team)]}"
        res.level = 0 
    
        res = bp2.text_frame.add_paragraph()
        res.text = f"Bollinnehav för/emot: \n\t{self.stats.prints['possession'][self.stats.main_team]} / {self.stats.prints['possession'][self.stats.opposite_team(self.stats.main_team)]}"
        res.level = 0 


        #res = bp2.text_frame.add_paragraph()
        #res.text = f"Skottsnitt FÖR på mål/skottförsök: {self.stats.prints['shots on goal'][self.stats.main_team]} / {sum(self.stats.prints['shot types'][self.stats.main_team].values())}"
        #res.level = 0

    def make_comparative_report_shots_for_page(self) -> None:
        ''''makes the shots for page for comparitive report'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_color(slide)
        self.add_main_team_logo(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Skott- och \nmålstatistik FÖR'
        table_header = ['Skottyp', 'Antal mål', 'Antal skott', 'Mål- procent']
        x, y, cx, cy = Inches(0), Inches(2.3), Inches(4.9), Inches(1.5)
        # display in this order (hopefully best first)
        st_display_order = ['friställande', 'retur', 'centralt', 'inlägg', 'dribbling', 'fast', 'utifrån']
        
        # vänster sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[1] # first 1, then 3  
        bp1.text = f"{gf.get_nickname(self.other.main_team,'short')} {self.other.number_of_games} halvlekar"

        table_frame = slide.shapes.add_table(9, 4, x + Inches(5.1)*0, y, cx, cy)
        table = table_frame.table
        for j, text in enumerate(table_header):
            table.cell(0, j).text = text
        for j, st in enumerate(st_display_order):
            table.cell(j + 1, 0).text = st.title()
            table.cell(j + 1, 1).text = f"{round(self.other.prints['goal types'][self.other.main_team][st] / self.other.number_of_games, 1)}"
            table.cell(j + 1, 2).text = f"{round(self.other.prints['shot types'][self.other.main_team][st] / self.other.number_of_games, 1)}"
            table.cell(j + 1, 3).text = f"{round(0 if self.other.prints['shot types'][self.other.main_team][st] == 0 else self.other.prints['goal types'][self.other.main_team][st]/self.other.prints['shot types'][self.other.main_team][st] * 100, 1)} %"
        table.cell(8, 0).text = 'Totalt'
        table.cell(8, 1).text = f"{round(sum(self.other.prints['goal types'][self.other.main_team].values()) / self.other.number_of_games, 1)}"
        table.cell(8, 2).text = f"{round(sum(self.other.prints['shot types'][self.other.main_team].values())/ self.other.number_of_games, 1)}"
        table.cell(8, 3).text = f"{round(sum(self.other.prints['goal types'][self.other.main_team].values())/sum(self.other.prints['shot types'][self.other.main_team].values()) * 100, 1)} %"

        # höger sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[3] # first 1, then 3  
        bp1.text = ' - '.join(gf.get_nickname(team,'short') for team in self.ordered_teams)

        table_frame = slide.shapes.add_table(9, 4, x + Inches(5.1)*1, y, cx, cy)
        table = table_frame.table
        for j, text in enumerate(table_header):
            table.cell(0, j).text = text
        for j, st in enumerate(st_display_order):
            table.cell(j + 1, 0).text = st.title()
            # skulle naturligtvis ha gjort så här hela tiden, vad fan händer?!
            goals = self.stats.prints['goal types'][self.stats.main_team][st] if st in self.stats.prints['goal types'][self.stats.main_team] else 0
            shots = self.stats.prints['shot types'][self.stats.main_team][st] if st in self.stats.prints['shot types'][self.stats.main_team] else 0
            table.cell(j + 1, 1).text = f"{goals}"
            table.cell(j + 1, 2).text = f"{shots}"
            table.cell(j + 1, 3).text = f"{round(0 if shots == 0 else goals/shots * 100, 1) } %"
        table.cell(8, 0).text = 'Totalt'
        table.cell(8, 1).text = f"{sum(self.stats.prints['score'][self.stats.main_team].values())}"
        table.cell(8, 2).text = f"{sum(self.stats.prints['shot types'][self.stats.main_team].values())}"
        table.cell(8, 3).text = f"{round(sum(self.stats.prints['score'][self.stats.main_team].values())/sum(self.stats.prints['shot types'][self.stats.main_team].values()) * 100, 1)} %"
        
        # remove the two text boxes that arent used
        for x in [2, 4]:
            subtitle = slide.placeholders[x]
            sp = subtitle.element
            sp.getparent().remove(sp)   

    def make_comparative_report_shots_against_page(self) -> None:
        ''''makes the shots for page against comparitive report'''
        ''''makes the shots for page for comparitive report'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_color(slide)
        #self.add_main_team_logo(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Skott- och \nmålstatistik EMOT'
        table_header = ['Skottyp', 'Antal mål', 'Antal skott', 'Mål- procent']
        x, y, cx, cy = Inches(0), Inches(2.3), Inches(4.9), Inches(1.5)
        # display in this order (hopefully best first)
        st_display_order = ['friställande', 'retur', 'centralt', 'inlägg', 'dribbling', 'fast', 'utifrån']
        
        # vänster sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[1] # first 1, then 3  
        bp1.text = f"{gf.get_nickname(self.other.main_team,'short')} {self.other.number_of_games} halvlekar"

        table_frame = slide.shapes.add_table(9, 4, x + Inches(5.1)*0, y, cx, cy)
        table = table_frame.table
        for j, text in enumerate(table_header):
            table.cell(0, j).text = text
        for j, st in enumerate(st_display_order):
            table.cell(j + 1, 0).text = st.title()
            table.cell(j + 1, 1).text = f"{round(self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)][st] / self.other.number_of_games, 1)}"
            table.cell(j + 1, 2).text = f"{round(self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)][st] / self.other.number_of_games, 1)}"
            table.cell(j + 1, 3).text = f"{round(0 if self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)][st] == 0 else self.other.prints['goal types'][self.other.main_team][st]/self.other.prints['shot types'][self.other.main_team][st] * 100, 1)} %"
        table.cell(8, 0).text = 'Totalt'
        table.cell(8, 1).text = f"{round(sum(self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)].values()) / self.other.number_of_games, 1)}"
        table.cell(8, 2).text = f"{round(sum(self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)].values())/ self.other.number_of_games, 1)}"
        table.cell(8, 3).text = f"{round(sum(self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)].values())/sum(self.other.prints['shot types'][self.other.main_team].values()) * 100, 1)} %"

        # höger sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[3] # first 1, then 3  
        bp1.text = ' - '.join(gf.get_nickname(team,'short') for team in self.ordered_teams)

        table_frame = slide.shapes.add_table(9, 4, x + Inches(5.1)*1, y, cx, cy)
        table = table_frame.table
        for j, text in enumerate(table_header):
            table.cell(0, j).text = text
        for j, st in enumerate(st_display_order):
            table.cell(j + 1, 0).text = st.title()
            table.cell(j + 1, 1).text = f"{self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)][st] if st in self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)] else 0}"
            table.cell(j + 1, 2).text = f"{self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)][st] if st in self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)] else 0}"
            table.cell(j + 1, 3).text = f"{round(self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)][st]/self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)][st] * 100 if st in self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)] else 0, 1)} %"
        table.cell(8, 0).text = 'Totalt'
        table.cell(8, 1).text = f"{sum(self.stats.prints['score'][self.stats.opposite_team(self.stats.main_team)].values())}"
        table.cell(8, 2).text = f"{sum(self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)].values())}"
        table.cell(8, 3).text = f"{round(sum(self.stats.prints['score'][self.stats.opposite_team(self.stats.main_team)].values())/sum(self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)].values()) * 100, 1)} %"
        
        # remove the two text boxes that arent used
        for x in [2, 4]:
            subtitle = slide.placeholders[x]
            sp = subtitle.element
            sp.getparent().remove(sp)   


    def make_season_report_overview_stats_page(self) -> None:
        '''makes the overview stats page layout'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Övergripande \nsäsongsstatistik'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Totala mål: \n\t{sum(self.stats.prints['goal types'][team].values())} - {round(sum(self.stats.prints['goal types'][team].values())/(sum(self.stats.prints['goal types'][team].values()) + sum(self.stats.prints['goal types'][self.stats.opposite_team(team)].values()))* 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Skott på mål (skottförsök): \n\t{self.stats.prints['shots on goal'][team]} ({self.stats.prints['shot attempts'][team]}) - {round(self.stats.prints['shots on goal'][team]/(self.stats.prints['shots on goal'][team] + self.stats.prints['shots on goal'][self.stats.opposite_team(team)]) * 100)} % ({round(self.stats.prints['shot attempts'][team]/(self.stats.prints['shot attempts'][team] + self.stats.prints['shot attempts'][self.stats.opposite_team(team)]) * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Närkamper och brytningar: \n\t{self.stats.prints['scrimmages'][team] + self.stats.prints['interceptions'][team]} - {round((self.stats.prints['scrimmages'][team] + self.stats.prints['interceptions'][team])/(self.stats.prints['scrimmages'][team] + self.stats.prints['interceptions'][team] + self.stats.prints['scrimmages'][self.stats.opposite_team(team)] + self.stats.prints['interceptions'][self.stats.opposite_team(team)]) * 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph() 
            res.text = f"Hörnornmål (hörnor): \n\t{self.stats.prints['score'][team]['hörnmål']} ({sum(self.stats.prints['corners'][team].values())}) - {round(self.stats.prints['score'][team]['hörnmål'] / (self.stats.prints['score'][team]['hörnmål'] + self.stats.prints['score'][self.stats.opposite_team(team)]['hörnmål'] )*100)} % ({round(sum(self.stats.prints['corners'][team].values()) / (sum(self.stats.prints['corners'][team].values()) + sum(self.stats.prints['corners'][self.stats.opposite_team(team)].values()) )*100)} %) "
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Bollinnehav: \n\t{self.stats.prints['possession'][team]} - {round(gf.readable_to_sec(self.stats.prints['possession'][team])/(gf.readable_to_sec(self.stats.prints['possession'][team]) + gf.readable_to_sec(self.stats.prints['possession'][self.stats.opposite_team(team)]))* 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            #res = bp2.text_frame.add_paragraph() # orkar inte göra turnary för noll
            #res.text = f"Utvisningar: \n\t{self.stats.prints['penalties'][team]}" # - {round(gf.readable_to_sec(self.stats.prints['possession'][team])/(gf.readable_to_sec(self.stats.prints['possession'][team]) + gf.readable_to_sec(self.stats.prints['possession'][self.stats.opposite_team(team)]))* 100)} %"
            #res.level = 0
            #res.font.color.rgb = self.get_team_text_color(team) 
  
    def make_season_report_shots_page(self, show_xg: bool) -> None:
        '''makes the shot stats page layout'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        #self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Skott- och \nmålstatistik'
        # list of shot types and total goals scored
        st_and_goals = [(st, sum(self.stats.prints['goal types'][team][st] for team in self.ordered_teams)) for st in self.stats.prints['goal types'][self.ordered_teams[0]]]
        # display the shot types in order of total scored goals
        st_display_order = [x[0] for x in sorted(st_and_goals, key = lambda x : x[1], reverse = True)]
        table_header = ['Skottyp', 'Antal mål', 'Antal skott', 'Mål- procent']
        x, y, cx, cy = Inches(0), Inches(2.3), Inches(4.9), Inches(1.5)
        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            xg_string = f", XG: {round(self.stats.prints['expected goals'][team], 2)}"
            bp1.text = f"{gf.get_nickname(team,'short')}{xg_string * show_xg}"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            table_frame = slide.shapes.add_table(9, 4, x + Inches(5.1)*i, y, cx, cy)
            table = table_frame.table
            for j, text in enumerate(table_header):
                table.cell(0, j).text = text
            for j, st in enumerate(st_display_order):
                table.cell(j + 1, 0).text = st.title()
                table.cell(j + 1, 1).text = f"{self.stats.prints['goal types'][team][st]}"
                table.cell(j + 1, 2).text = f"{self.stats.prints['shot types'][team][st]}"
                table.cell(j + 1, 3).text = f"{round(0 if self.stats.prints['shot types'][team][st] == 0 else self.stats.prints['goal types'][team][st]/self.stats.prints['shot types'][team][st] * 100, 1)} %"
            table.cell(8, 0).text = 'Totalt'
            table.cell(8, 1).text = f"{sum(self.stats.prints['goal types'][team].values())}"
            table.cell(8, 2).text = f"{sum(self.stats.prints['shot types'][team].values())}"
            table.cell(8, 3).text = f"{round(sum(self.stats.prints['goal types'][team].values())/sum(self.stats.prints['shot types'][team].values()) * 100, 1)} %"
        
           # for cell in PP.iter_cells(table):
           #     cell.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 
            #    cell.fill.solid()
             #   cell.fill.fore_color.rgb = PP.background_color

        # remove the two text boxes that arent used
        for x in [2, 4]:
            subtitle = slide.placeholders[x]
            sp = subtitle.element
            sp.getparent().remove(sp)

    def make_season_report_duels_page(self) -> None:
        '''makes the duels stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Bollvinster'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Vunna närkamper: \n\t{self.stats.prints['scrimmages'][team]} ({round(self.stats.prints['scrimmages'][team] / (self.stats.prints['scrimmages'][team] + self.stats.prints['scrimmages'][self.stats.opposite_team(team)]) * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team)
            res = bp2.text_frame.add_paragraph()
            res.text = f"Brytningar: \n\t{self.stats.prints['interceptions'][team]} ({round(self.stats.prints['interceptions'][team] / (self.stats.prints['interceptions'][team] + self.stats.prints['interceptions'][self.stats.opposite_team(team)]) * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Bolltapp: \n\t{self.stats.prints['lost balls'][team]} ({round(self.stats.prints['lost balls'][team] / (self.stats.prints['lost balls'][team] + self.stats.prints['lost balls'][self.stats.opposite_team(team)]) * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
    
    def make_season_report_long_shot_targets_page(self) -> None:
        '''makes the targets of long shots stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Resultat av \nskott utifrån'

        for i, team in enumerate(self.ordered_teams):
            total_shots = sum(self.stats.prints['long shots outcomes'][team].values())
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = f"{gf.get_nickname(team,'short')}: {total_shots} st"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Mål: \n\t{self.stats.prints['long shots outcomes'][team]['mål']} ({round(self.stats.prints['long shots outcomes'][team]['mål']/total_shots * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Utanför: \n\t{self.stats.prints['long shots outcomes'][team]['utanför']} ({round(self.stats.prints['long shots outcomes'][team]['utanför']/total_shots * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Täckta: \n\t{self.stats.prints['long shots outcomes'][team]['täckt']} ({round(self.stats.prints['long shots outcomes'][team]['täckt']/total_shots * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Räddningar: \n\t{self.stats.prints['long shots outcomes'][team]['räddning']} ({round(self.stats.prints['long shots outcomes'][team]['räddning']/total_shots * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
    
    def make_season_report_possession_after_long_shots_page(self) -> None:
        '''makes the possession after long shots stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Bollinnehav efter \nskott utifrån'

        for i, team in enumerate(self.ordered_teams):
            total_shots = sum(self.stats.prints['after long shots per teams'][team].values())
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = f"{gf.get_nickname(team,'short')}: {total_shots} st"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            opponent = self.stats.opposite_team(team)
            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            #res.text = f"Mål: \n\t{self.stats.prints['long shots outcomes'][team]['mål']} ({round(self.stats.prints['long shots outcomes'][team]['mål']/total_shots * 100)} %)"
            res.text = f"{gf.get_nickname(team,'short')}: \n\t{self.stats.prints['after long shots per teams'][team][team]} ({round(self.stats.prints['after long shots per teams'][team][team]/total_shots * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"{gf.get_nickname(opponent,'short')}: \n\t{self.stats.prints['after long shots per teams'][team][opponent]} ({round(self.stats.prints['after long shots per teams'][team][opponent]/total_shots * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 

    def make_season_report_after_long_shots_attacking_team_page(self) -> None:
        '''makes the attacking team long shots stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Händelsen efter skott\nutifrån: anfallare'
        after = self.stats.prints['after long shots events']

        for i, team in enumerate(self.ordered_teams):
            attacker = after[team][team]
            total_times = sum(self.stats.prints['after long shots per teams'][team].values())
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = f"{gf.get_nickname(team,'short')}: {total_times} st"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            #opponent = self.stats.opposite_team(team)
            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            goals = attacker['mål'] if 'mål' in attacker else 0
            res.text = f"Mål: \n\t{goals} ({round(goals/total_times * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            shots = attacker['skott'] if 'skott' in attacker else 0
            res.text = f"Skott: \n\t{shots + goals} ({round((shots + goals)/total_times * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            corner = attacker['hörna'] if 'skott' in attacker else 0
            res.text = f"Hörnor: \n\t{corner} ({round(corner/total_times * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            #res = bp2.text_frame.add_paragraph()
            other = total_times - goals - corner - shots
            #res.text = f"Övrigt: \n\t{other} ({round(other/total_times * 100)} %)"
            #res.level = 0
            #res.font.color.rgb = self.get_team_text_color(team) 

    def make_season_report_after_long_shots_defending_team_page(self) -> None:
        '''makes the defending team long shots stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Händelsen efter skott\nutifrån: försvarare'
        after = self.stats.prints['after long shots events']

        for i, team in enumerate(self.ordered_teams):
            defender = after[self.stats.opposite_team(team)][team]
            total_times = sum(self.stats.prints['after long shots per teams'][self.stats.opposite_team(team)].values())
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = f"{gf.get_nickname(team,'short')}: {total_times} st"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            #opponent = self.stats.opposite_team(team)
            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            throw = defender['utkast'] if 'utkast' in defender else 0
            res.text = f"Utkast: \n\t{throw} ({round(throw/total_times * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            clean = defender['rensning'] if 'rensning' in defender else 0 # kul att kalla det clean
            res.text = f"Rensning: \n\t{clean} ({round((clean)/total_times * 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            #res = bp2.text_frame.add_paragraph()
            #corner = defender['hörna'] if 'skott' in defender else 0
            #res.text = f"Hörnor: \n\t{corner} ({round(corner/total_times * 100)} %)"
            #res.level = 0
            #res.font.color.rgb = self.get_team_text_color(team) 
            #res = bp2.text_frame.add_paragraph()
            #other = total_times - goals - corner - shots
            #res.text = f"Övrigt: \n\t{other} ({round(other/total_times * 100)} %)"
            #res.level = 0
            #res.font.color.rgb = self.get_team_text_color(team) 


    def make_comparative_report_slot_for_page(self) -> None:
        ''''makes the offensive slots page for the comparison report'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_main_team_logo(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Offensiva inspel'

        # vänster sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[1] # first 1, then 3  
        bp1.text = f"{gf.get_nickname(self.other.main_team,'short')} {self.other.number_of_games} halvlekar"
        #bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

        bp2 = bpb.placeholders[2]  # first 2, then 4
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsmål/inspelsskott: \n\t{round(self.other.prints['goal types'][self.other.main_team]['inlägg'] / self.other.number_of_games, 1) if 'inlägg' in self.other.prints['goal types'][self.other.main_team] else 0} / {round(self.other.prints['shot types'][self.other.main_team]['inlägg'] / self.other.number_of_games, 1) if 'inlägg' in self.other.prints['shot types'][self.other.main_team] else 0} = {round(self.other.prints['goal types'][self.other.main_team]['inlägg'] /self.other.prints['shot types'][self.other.main_team]['inlägg'], 1) * 100 if 'inlägg' in self.other.prints['goal types'][self.other.main_team] else 0} %" 
        res.level = 0
        #res.font.color.rgb = self.get_team_text_color(team) 
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsskott/inspel: \n\t{round(self.other.prints['shot types'][self.other.main_team]['inlägg'] / self.other.number_of_games, 1) if 'inlägg' in self.other.prints['shot types'][self.other.main_team] else 0} / {round(self.other.prints['slot passes'][self.other.main_team] / self.other.number_of_games, 1)} = {round(self.other.prints['shot types'][self.other.main_team]['inlägg'] /self.other.prints['slot passes'][self.other.main_team], 1) * 100 if 'inlägg' in self.other.prints['goal types'][self.other.main_team] else 0} %" 
        res.level = 0 

        # höger sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[3] # first 1, then 3  
        bp1.text = ' - '.join(gf.get_nickname(team,'short') for team in self.ordered_teams)
        #bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

        bp2 = bpb.placeholders[4]  # first 2, then 4
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsmål/inspelsskott: \n\t{self.stats.prints['goal types'][self.stats.main_team]['inlägg'] if 'inlägg' in self.stats.prints['goal types'][self.stats.main_team] else 0} / {self.stats.prints['shot types'][self.stats.main_team]['inlägg'] if 'inlägg' in self.other.prints['shot types'][self.stats.main_team] else 0} = {round(self.stats.prints['goal types'][self.stats.main_team]['inlägg'] /self.stats.prints['shot types'][self.stats.main_team]['inlägg'], 1) * 100 if 'inlägg' in self.stats.prints['goal types'][self.stats.main_team] else 0} %" 
        res.level = 0
        #res.font.color.rgb = self.get_team_text_color(team) 
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsskott/inspel: \n\t{self.stats.prints['shot types'][self.stats.main_team]['inlägg'] if 'inlägg' in self.stats.prints['shot types'][self.stats.main_team] else 0} / {self.stats.prints['slot passes'][self.stats.main_team]} = {round(self.stats.prints['shot types'][self.stats.main_team]['inlägg'] /self.stats.prints['slot passes'][self.stats.main_team], 1) * 100 if 'inlägg' in self.stats.prints['goal types'][self.stats.main_team] else 0} %" 
        res.level = 0 
    
    def make_comparative_report_slot_against_page(self) -> None:
        '''makes the defensive slots page for the comparison report'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        #self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Defensiva inspel'

        # vänster sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[1] # first 1, then 3  
        bp1.text = f"{gf.get_nickname(self.other.main_team,'short')} {self.other.number_of_games} halvlekar"
        #bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

        bp2 = bpb.placeholders[2]  # first 2, then 4
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsmål/inspelsskott: \n\t{round(self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)]['inlägg'] if 'inlägg' in self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)] else 0 / self.other.number_of_games, 1)} / {round(self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)]['inlägg'] if 'inlägg' in self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)] else 0 / self.other.number_of_games, 1)} = {round(self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)]['inlägg'] /self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)]['inlägg'], 1) * 100 if 'inlägg' in self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)] else 0} %" 
        res.level = 0
        #res.font.color.rgb = self.get_team_text_color(team) 
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsskott/inspel: \n\t{round(self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)]['inlägg'] if 'inlägg' in self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)] else 0 / self.other.number_of_games, 1)} / {round(self.other.prints['slot passes'][self.other.opposite_team(self.other.main_team)] / self.other.number_of_games, 1)} = {round(self.other.prints['shot types'][self.other.opposite_team(self.other.main_team)]['inlägg'] /self.other.prints['slot passes'][self.other.opposite_team(self.other.main_team)], 1) * 100 if 'inlägg' in self.other.prints['goal types'][self.other.opposite_team(self.other.main_team)] else 0} %" 
        res.level = 0 
        res.level = 0

        # höger sida
        bpb = slide.shapes
        bp1 = bpb.placeholders[3] # first 1, then 3  
        bp1.text = ' - '.join(gf.get_nickname(team,'short') for team in self.ordered_teams)
        #bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

        bp2 = bpb.placeholders[4]  # first 2, then 4
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsmål/inspelsskott: \n\t{self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)]['inlägg'] if 'inlägg' in self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)] else 0} / {self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)]['inlägg'] if 'inlägg' in self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)] else 0} = {round(self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)]['inlägg'] /self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)]['inlägg'], 1) * 100 if 'inlägg' in self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)] else 0} %" 
        res.level = 0
        #res.font.color.rgb = self.get_team_text_color(team)  
        res = bp2.text_frame.add_paragraph()
        res.text = f"Inspelsskott/inspel: \n\t{self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)]['inlägg'] if 'inlägg' in self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)] else 0} / {self.stats.prints['slot passes'][self.stats.opposite_team(self.stats.main_team)]}  = {round(self.stats.prints['shot types'][self.stats.opposite_team(self.stats.main_team)]['inlägg'] /self.stats.prints['slot passes'][self.stats.opposite_team(self.stats.main_team)], 1) * 100 if 'inlägg' in self.stats.prints['goal types'][self.stats.opposite_team(self.stats.main_team)] else 0} %" 
        res.level = 0 
        res.level = 0


    def make_game_report_slot_page(self) -> None:
        '''makes the slot passes stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Inspel'
        slot_goals = {t: len([x for x in self.stats.goals_info_list if x['team'] == t and x['shot type'] == 'inlägg']) for t in self.ordered_teams}

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Inspelsmål/inspelsskott: \n\t{slot_goals[team]}/{0 if 'inlägg' not in self.stats.prints['shot types'][team] else self.stats.prints['shot types'][team]['inlägg']} = {round(0 if 'inlägg' not in self.stats.prints['shot types'][team] else slot_goals[team] / self.stats.prints['shot types'][team]['inlägg'] * 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team)  
            res = bp2.text_frame.add_paragraph()
            
            inspel = inspel = self.stats.prints['slot passes'][team]
            if inspel == 0:
                skott = 0
                andel = 0
            else:
                skott = 0 if 'inlägg' not in self.stats.prints['shot types'][team] else self.stats.prints['shot types'][team]['inlägg']
                andel = round(skott/inspel * 100, 1)
            #res.text = f"Inspelsskott/inspel: \n\t{0 if 'inlägg' not in self.stats.prints['shot types'][team] else self.stats.prints['shot types'][team]['inlägg']}/{self.stats.prints['slot passes'][team]} = {round(0 if 'inlägg' not in self.stats.prints['shot types'][team] else self.stats.prints['shot types'][team]['inlägg'] / self.stats.prints['slot passes'][team] * 100)} %"
            res.text = f"Inspelsskott/inspel: \n\t{skott}/{inspel}={andel} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            #res = bp2.text_frame.add_paragraph()
            #res.text = f"Totala % mål, skott, inspel: \n\t{round(self.stats.prints['goal types'][team]['inlägg']/(self.stats.prints['goal types'][team]['inlägg'] + self.stats.prints['goal types'][self.stats.opposite_team(team)]['inlägg'])* 100, 1)} %,  {round(self.stats.prints['shot types'][team]['inlägg']/(self.stats.prints['shot types'][team]['inlägg'] + self.stats.prints['shot types'][self.stats.opposite_team(team)]['inlägg'])* 100, 1)} %, {round(self.stats.prints['slot passes'][team]/(self.stats.prints['slot passes'][team] + self.stats.prints['slot passes'][self.stats.opposite_team(team)]) * 100, 1)} %"
            #res.level = 0
            #res.font.color.rgb = self.get_team_text_color(team) 


    def make_season_report_slot_page(self) -> None:
        '''makes the slot passes stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Inspel'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Inspelsmål/inspelsskott: \n\t{self.stats.prints['goal types'][team]['inlägg']}/{self.stats.prints['shot types'][team]['inlägg']} = {round(0 if self.stats.prints['shot types'][team]['inlägg'] == 0 else self.stats.prints['goal types'][team]['inlägg'] / self.stats.prints['shot types'][team]['inlägg'] * 100, 1)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Inspelsskott/inspel: \n\t{self.stats.prints['shot types'][team]['inlägg']}/{self.stats.prints['slot passes'][team]} = {round(0 if self.stats.prints['slot passes'][team] == 0 else self.stats.prints['shot types'][team]['inlägg'] / self.stats.prints['slot passes'][team] * 100, 1)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            try:
                goal_percentage = round(self.stats.prints['goal types'][team]['inlägg']/(self.stats.prints['goal types'][team]['inlägg'] + self.stats.prints['goal types'][self.stats.opposite_team(team)]['inlägg'])* 100, 1)
            except ZeroDivisionError:
                goal_percentage = 0
            try:
                shot_percentage = round(self.stats.prints['shot types'][team]['inlägg']/(self.stats.prints['shot types'][team]['inlägg'] + self.stats.prints['shot types'][self.stats.opposite_team(team)]['inlägg'])* 100, 1)
            except ZeroDivisionError:
                shot_percentage = 0
            try:
                slot_pass_percentage = round(self.stats.prints['slot passes'][team]/(self.stats.prints['slot passes'][team] + self.stats.prints['slot passes'][self.stats.opposite_team(team)]) * 100, 1)
            except ZeroDivisionError:
                slot_pass_percentage = 0
            #res.text = f"Totala % mål, skott, inspel: \n\t{0 if 'inlägg' not in self.stats.prints['goal types'][team]['inlägg'] else round(self.stats.prints['goal types'][team]['inlägg']/(self.stats.prints['goal types'][team]['inlägg'] + self.stats.prints['goal types'][self.stats.opposite_team(team)]['inlägg'])* 100, 1)} %,  {round(self.stats.prints['shot types'][team]['inlägg']/(self.stats.prints['shot types'][team]['inlägg'] + self.stats.prints['shot types'][self.stats.opposite_team(team)]['inlägg'])* 100, 1)} %, {round(self.stats.prints['slot passes'][team]/(self.stats.prints['slot passes'][team] + self.stats.prints['slot passes'][self.stats.opposite_team(team)]) * 100, 1)} %"
            res.text = f"Totala % mål, skott, inspel: \n\t{goal_percentage} % {shot_percentage} % {slot_pass_percentage} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 

    def make_season_report_corners_page(self) -> None:
        '''makes the corners stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Hörnor'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Högermål/högerhörnor: \n\t{self.stats.prints['corner goal sides'][team]['right']}/{self.stats.prints['corners'][team]['right']} = {round(0 if self.stats.prints['corners'][team]['right'] == 0 else self.stats.prints['corner goal sides'][team]['right'] / self.stats.prints['corners'][team]['right'] * 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Vänstermål/vänsterhörnor: \n\t{self.stats.prints['corner goal sides'][team]['left']}/{self.stats.prints['corners'][team]['left']} = {round(0 if self.stats.prints['corners'][team]['left'] == 0 else self.stats.prints['corner goal sides'][team]['left'] / self.stats.prints['corners'][team]['left'] * 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Mål/totala hörnor: \n\t{self.stats.prints['corner goal sides'][team]['left'] + self.stats.prints['corner goal sides'][team]['right']}/{self.stats.prints['corners'][team]['right'] + self.stats.prints['corners'][team]['left']} = {round(0 if (self.stats.prints['corners'][team]['left'] + self.stats.prints['corners'][team]['right']) == 0 else (self.stats.prints['corner goal sides'][team]['left'] + self.stats.prints['corner goal sides'][team]['right'])/ (self.stats.prints['corners'][team]['left'] + self.stats.prints['corners'][team]['right']) * 100)} %"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 


    def make_game_report_overview_stats_page(self) -> None:
        '''makes the overview stats page layout'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Matchstatistik'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Resultat: \n\t{sum(self.stats.prints['score'][team].values())}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"XG: \n\t{round(self.stats.prints['expected goals'][team],2)}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Skott på mål: \n\t{self.stats.prints['shots on goal'][team]}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            corner_goals = self.stats.prints['score'][team]['hörnmål'] if 'hörnmål' in self.stats.prints['score'][team] else 0
            res.text = f"Hörnor (mål): \n\t{sum(self.stats.prints['corners'][team].values())} ({corner_goals})"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team)
            res = bp2.text_frame.add_paragraph()
            res.text = f"Bollinnehav: \n\t{self.stats.prints['possession'][team]} ({round(gf.readable_to_sec(self.stats.prints['possession'][team])/(gf.readable_to_sec(self.stats.prints['possession'][team]) + gf.readable_to_sec(self.stats.prints['possession'][self.stats.opposite_team(team)]))* 100)} %)"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 

    def make_game_report_shot_types_page(self, show_xg: bool) -> None:
        '''makes the shots types stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Skottstatistik \nSkotttyper'
        # this is just all the shot types for both teams sorted in alphabetical order
        print_order = sorted(list({st for team in self.stats.prints['shot types'] for st in self.stats.prints['shot types'][team]}))
        for i, team in enumerate(self.ordered_teams):
            xg_string = f", XG: {round(self.stats.prints['expected goals'][team], 2)}"
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = f"{gf.get_nickname(team,'short')} - skott: {sum(self.stats.prints['shot types'][team].values())}{xg_string * show_xg}"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            for st in print_order:
                res = bp2.text_frame.add_paragraph()
                if st in self.stats.prints['shot types'][team]:
                    i = self.stats.prints['shot types'][team][st]
                else: 
                    i = 0
                res.text = f"{st.title()}: {i}"
                res.level = 0
                res.font.color.rgb = self.get_team_text_color(team) 

    def make_game_report_duels_page(self) -> None:
        '''makes the duels stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Bollvinster'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Vunna närkamper: \n\t{self.stats.prints['scrimmages'][team]}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Brytningar: \n\t{self.stats.prints['interceptions'][team]}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Bolltapp: \n\t{self.stats.prints['lost balls'][team]}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 

    def make_game_report_scimmages_page(self) -> None:
        '''makes the scrimmages stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Närkampssituationer \noch deras utfall'

        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = gf.get_nickname(team,'short')
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 
            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            res = bp2.text_frame.add_paragraph()
            res.text = f"Långvarig bollvinst: \n\t{self.stats.prints['possession changes'][team]['long']}"
            res.level = 0
            res.font.color.rgb = self.get_team_text_color(team) 
            res = bp2.text_frame.add_paragraph()
            res.text = f"Kortvarig bollvinst: \n\t{self.stats.prints['possession changes'][team]['short']}"
            res.font.color.rgb = self.get_team_text_color(team) 
            res.level = 0

    def make_game_report_shot_origins_page(self) -> None:
        '''makes the shots types stats page'''
        slide_register = self.pres.slide_layouts[4]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        title.text = 'Skottstatistik \nSkottens ursprung'
        # this is the shot origins for both teams sorted 
        print_order = sorted(list({so for team in self.stats.prints['shot origins'] for so in self.stats.prints['shot origins'][team]}))


        for i, team in enumerate(self.ordered_teams):
            bpb = slide.shapes
            bp1 = bpb.placeholders[(i + 1)*2-1] # first 1, then 3  
            bp1.text = f"{gf.get_nickname(team,'short')} - skottförsök: {sum(self.stats.prints['shot origins'][team].values())}"
            bp1.text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            bp2 = bpb.placeholders[(i + 1)*2]  # first 2, then 4
            #bp2.text_frame.paragraphs[0].font.color.rgb = constants.colors[team][0]
            for so in print_order:
                res = bp2.text_frame.add_paragraph()
                if so in self.stats.prints['shot origins'][team]:
                    i = self.stats.prints['shot origins'][team][so]
                else: 
                    i = 0
                res.text = f"{so.title()}: {i}"
                res.level = 0
                res.font.color.rgb = self.get_team_text_color(team) 

    def make_game_report_goals_stats_page(self) -> None:
        '''makes the goals stats page'''
        slide_register = self.pres.slide_layouts[1]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide)
        title = slide.shapes.title
        title.text = 'Målen'
        bpb = slide.shapes

        bp1 = bpb.placeholders[1]
        for goal in self.stats.goals_info_list:
            res = bp1.text_frame.add_paragraph()
            # fixa det här med målens tid i andra halvlek
            res.text = f"{goal['time']}: {goal['origin']} -> {goal['shot type']} på {goal['attack time']}s."
            res.font.color.rgb = self.get_team_text_color(goal['team']) 
            res.level = 0

    def make_game_report_before_and_after_table_page(self) -> None:
        '''makes the page with the table for the before and after possessions'''
        slide_register = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title

        title.text = 'Närkamper \noch deras utfall'
        left_table = Inches(0.25)
        top_table = Inches(2.5)
        width_table = Inches(5)
        height_table = Inches(4)
        table_frame = slide.shapes.add_table(3,3, left_table, top_table, width_table, height_table)
        table = table_frame.table
        for i, team in enumerate(self.stats.prints['before and after']):
            table.cell(i+1, 0).text = f"Före {gf.get_nickname(team,'short')}"
            table.cell(i+1, 0).text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            table.cell(0, i+1).text = f"Efter {gf.get_nickname(team,'short')}"
            table.cell(0, i+1).text_frame.paragraphs[0].font.color.rgb = self.get_team_text_color(team) 

            table.cell(i+1, i+1).text = f"{self.stats.prints['before and after'][team][team]}"
            table.cell(i+1, i+1).fill.solid()
            # background is team's secondary color with opacity 80 % against a white background
            table.cell(i+1, i+1).fill.fore_color.rgb = PP.background_color#RGBColor(*gf.faded_rgb_color(constants.colors[team][1], 0.8))
            if i == 0:
                j,k = 1,2
            else:
                j,k = 2,1
            # a blend of the two team's colors at 50 %
            blended_color = RGBColor(*gf.faded_rgb_color(constants.colors[self.stats.teams.copy().pop()][1], 0.5, background = constants.colors[self.stats.opposite_team(self.stats.teams.copy().pop())][1]))
            # fading the color at 60 % on white background
            faded_color = RGBColor(*gf.faded_rgb_color(blended_color, 0.6))
            table.cell(j,k).text = f"{self.stats.prints['before and after'][team][self.stats.opposite_team(team)]}"
            table.cell(j,k).fill.solid()
            table.cell(j,k).fill.fore_color.rgb = PP.background_color # faded_color
        self.style_before_and_after_table(table_frame)

        duels_list = [{team: 0 for team in self.stats.teams} for i in range(2)]
        for team in self.stats.prints['before and after'][self.stats.main_team]:
            duels_list[0][team] = self.stats.prints['before and after'][self.stats.main_team][team]
        for team in self.stats.prints['before and after'][self.stats.opposite_team(self.stats.main_team)]:
            duels_list[1][team] = self.stats.prints['before and after'][self.stats.opposite_team(self.stats.main_team)][team]
        x_labels = [f"Fördel {gf.get_nickname(self.stats.main_team,'full')}", f"Fördel {gf.get_nickname(self.stats.opposite_team(self.stats.main_team),'full')}"]
        image = self.plot.make_value_vertical_bars(duels_list, title='Närkampers utfall givet lags innehav före', ylabel='Närkamper (antal)', xlabel='Spelförande lag innan närkamp', x_labels = x_labels, main_team_color=gf.rgb255_to_rgb1(PP.image_color_main), other_team_color=gf.rgb255_to_rgb1(PP.image_color_opponent)) #other_team_color=gf.rgb255_to_rgb1(constants.colors[self.stats.opposite_team(self.stats.main_team)][1]))
        slide.shapes.add_picture(image, width_table, top_table, width_table + Inches(0.5))


    def make_single_image_page(self, image_link: str, title_text: str, from_left = 0, from_top = 0.5, width = 10) -> None:
        '''makes a page with one centered image and a title text'''
        slide_register = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(slide_register)
        self.set_background_color(slide)
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = PP.text_color

        img = image_link
        slide.shapes.add_picture(img, Inches(from_left), Inches(from_top), Inches(width))


    def make_game_report_attacks_and_fourty_page(self) -> None:
        '''makes a page with the 40 and sustained attacks data'''
        slide_register = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        #self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title

        title.text = '40-situationer och långa anfall'
        fourty_image = self.plot.make_per_minute_bars(self.stats.prints['40'], title='40-spel under matchen', ylabel='40-spel (antal)', color=gf.rgb255_to_rgb1(PP.image_color_main))
        sustained_attacks_image = self.plot.make_team_minute_bars(self.stats.prints['sustained attacks'], main_team_color=gf.rgb255_to_rgb1(PP.image_color_main), other_team_color= gf.rgb255_to_rgb1(PP.image_color_opponent), title= 'Långa anfall', ylabel='sekunder')

        from_left = Inches(0.25)
        from_top = Inches(2.5)
        width_image = Inches(5.5)
        height_image = Inches(4)
        
        slide.shapes.add_picture(fourty_image, Inches(0), from_top, width_image)#, height=height_image, width = width_image)
        slide.shapes.add_picture(sustained_attacks_image, width_image - Inches(0.6), from_top, width_image) #+ from_left * 2)#, height=height_image)

    def make_game_report_per_time_page(self) -> None:
        '''makes the page with the parts of game stats'''
        slide_register = self.pres.slide_layouts[5]
        slide = self.pres.slides.add_slide(slide_register)
        #self.set_background_image(slide)
        self.set_background_color(slide)
        #self.add_logo_images(slide, width = 1.5)
        title = slide.shapes.title
        main_team_color = gf.rgb255_to_rgb1(PP.image_color_main)
        opponent_color = gf.rgb255_to_rgb1(PP.image_color_opponent)
        title.text = f"Halvleken uppdelad i {gf.readable_number(self.stats.N)} delar"
        poss_image = self.plot.make_time_vertical_bars(self.stats.prints['per time lists']['possession'], xlabel='Halvleksdel', ylabel='Innehav (HH:MM:SS)', title='Bollinnehav', main_team_color = main_team_color, other_team_color= opponent_color)
        shots_image = self.plot.make_value_vertical_bars(self.stats.prints['per time lists']['shots'], xlabel='Halvleksdel', ylabel='Skottförsök (antal)', title='Skottförsök', main_team_color = main_team_color, other_team_color = opponent_color)
        duels_image = self.plot.make_value_vertical_bars(self.stats.prints['per time lists']['duels'], xlabel='Halvleksdel', ylabel='Närkamper och brytningar (antal)', title='Närkamper och brytningar', main_team_color = main_team_color, other_team_color = opponent_color)
        goals_image = self.plot.make_value_vertical_bars(self.stats.prints['per time lists']['goals'], xlabel='Halvleksdel', ylabel='Mål (antal)', title='Mål', main_team_color = main_team_color, other_team_color = opponent_color)
        slide.shapes.add_picture(poss_image, Inches(0.25), Inches(1), Inches(4.5))
        slide.shapes.add_picture(shots_image, Inches(5.25), Inches(1), Inches(4.5))
        slide.shapes.add_picture(duels_image, Inches(0.25), Inches(4.1), Inches(4.5))
        slide.shapes.add_picture(goals_image, Inches(5.25), Inches(4.1), Inches(4.5))

    def style_before_and_after_table(self, table_frame: pptx.shapes.graphfrm.GraphicFrame) -> None:
        '''sets the background and text colors of the table
            only used for the confusion matrix of before and after'''
        # table style
        tbl =  table_frame._element.graphic.graphicData.tbl
        style_id = '{69CF1AB2-1976-4502-BF36-3FF5EA218861}'
        tbl[0][-1].text = style_id
    
        table = table_frame.table
        # sets the outer cells to transparent 
        table.cell(0,0).fill.background()
        for i in range(1,3):
            table.cell(i,0).fill.background()
            table.cell(0,i).fill.background()
        
        # deals with text 
        for cell in PP.iter_cells(table):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(25)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER        

    def save_presentation(self) -> None:
        '''saves the presentation'''
        self.pres.save(self.stats.out + '.pptx')